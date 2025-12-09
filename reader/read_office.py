import os


def extensions_type_extract():
    """Return set of supported office document extensions"""
    return {
        '.docx',
        '.doc',
        '.xlsx',
        '.xls',
        '.csv',
        '.pptx',
        '.ppt'
    }


def specify_office_method_of_reading_the_file(file_info, logger=None):

    file_path = file_info.get("path")
    
    if not file_path or not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        # print(msg)
        # if logger:
        #     logger.error("Office file not found", file_path=file_path)
        return None
    
    file_path = str(file_path)
    file_lower = file_path.lower()
    
    try:
        if file_lower.endswith('.docx'):
            return read_docx_file(file_path)
        elif file_lower.endswith('.doc'):
            return read_doc_file(file_path, logger=logger)
        elif file_lower.endswith('.xlsx'):
            return read_xlsx_file(file_path)
        elif file_lower.endswith('.xls'):
            return read_xls_file(file_path)
        elif file_lower.endswith('.csv'):
            return read_csv_file(file_path)
        elif file_lower.endswith('.pptx'):
            return read_pptx_file(file_path)
        elif file_lower.endswith('.ppt'):
            print("⚠ .ppt format requires additional libraries")
            # print(msg)
            # if logger:
            #     logger.warning(msg, file_path=file_path)
            return None
        else:
            print(f"✗ Unsupported office file type: {file_path}")
            # print(msg)
            # if logger:
            #     logger.warning("Unsupported office file type", file_path=file_path)
            return None
    except Exception as e:
        print(f"✗ Error reading {file_path}: {str(e)}")
        # print(msg)
        # if logger:
        #     logger.error("Error reading office file", file_path=file_path, error=str(e))
        return None


def read_docx_file(filepath):
    """Independent DOCX reader function."""
    try:
        from docx import Document
    except ImportError:
        return {
            "error": "python-docx not installed. Install with: pip install python-docx",
            "filepath": filepath
        }
    
    try:
        if not os.path.exists(filepath):
            return {"error": "File not found", "filepath": filepath}
        
        doc = Document(filepath)
        
        result = {
            "filepath": filepath,
            "paragraphs": [],
            "tables": [],
            "total_paragraphs": 0,
            "total_tables": 0
        }
        
        for para in doc.paragraphs:
            if para.text.strip():
                result["paragraphs"].append({
                    "text": para.text,
                    "style": para.style.name if para.style else "Normal"
                })
        
        result["total_paragraphs"] = len(result["paragraphs"])
        
        for table_idx, table in enumerate(doc.tables):
            table_data = {
                "table_number": table_idx + 1,
                "rows": []
            }
            for row in table.rows:
                table_data["rows"].append([cell.text for cell in row.cells])
            result["tables"].append(table_data)
        
        result["total_tables"] = len(result["tables"])
        
        return result
    except Exception as e:
        return {"error": str(e), "filepath": filepath}


def read_doc_file(filepath, logger=None):
    """Independent DOC reader function - tries multiple methods."""
    if not os.path.exists(filepath):
        return {"error": "File not found", "filepath": filepath}
    
    # Try docx2txt
    try:
        import docx2txt
        # if logger:
        #     logger.debug("Attempting to read .doc with docx2txt", file_path=filepath)
        
        text = docx2txt.process(filepath)
        
        result = {
            "filepath": filepath,
            "text": text,
            "method": "docx2txt",
            "total_characters": len(text),
            "total_lines": len(text.splitlines())
        }
        
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if paragraphs:
            result["paragraphs"] = paragraphs
            result["total_paragraphs"] = len(paragraphs)
        
        return result
    except ImportError:
        pass
    except Exception as e:
        # if logger:
        #     logger.debug(f"docx2txt failed: {str(e)}", file_path=filepath)
        pass
    
    # Try antiword
    try:
        import subprocess
        import shutil
        
        antiword_path = shutil.which('antiword')
        if antiword_path:
            # if logger:
            #     logger.debug("Attempting to read .doc with antiword", file_path=filepath)
            
            result = subprocess.run(
                [antiword_path, filepath],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if result.returncode == 0:
                text = result.stdout
                
                doc_result = {
                    "filepath": filepath,
                    "text": text,
                    "method": "antiword",
                    "total_characters": len(text),
                    "total_lines": len(text.splitlines())
                }
                
                paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                if paragraphs:
                    doc_result["paragraphs"] = paragraphs
                    doc_result["total_paragraphs"] = len(paragraphs)
                
                return doc_result
    except Exception as e:
        # if logger:
        #     logger.debug(f"antiword failed: {str(e)}", file_path=filepath)
        pass
    
    # Try LibreOffice
    try:
        import subprocess
        import shutil
        import tempfile
        
        libreoffice_cmds = ['libreoffice', 'soffice']
        libreoffice_path = None
        
        for cmd in libreoffice_cmds:
            path = shutil.which(cmd)
            if path:
                libreoffice_path = path
                break
        
        if libreoffice_path:
            # if logger:
            #     logger.debug("Attempting to read .doc with LibreOffice", file_path=filepath)
            
            with tempfile.TemporaryDirectory() as tmpdir:
                try:
                    result = subprocess.run(
                        [
                            libreoffice_path,
                            '--headless',
                            '--convert-to', 'txt',
                            '--outdir', tmpdir,
                            filepath
                        ],
                        capture_output=True,
                        text=True,
                        timeout=60
                    )
                    
                    if result.returncode == 0:
                        base_name = os.path.splitext(os.path.basename(filepath))[0]
                        txt_file = os.path.join(tmpdir, f"{base_name}.txt")
                        
                        if os.path.exists(txt_file):
                            with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                                text = f.read()
                            
                            doc_result = {
                                "filepath": filepath,
                                "text": text,
                                "method": "libreoffice",
                                "total_characters": len(text),
                                "total_lines": len(text.splitlines())
                            }
                            
                            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                            if paragraphs:
                                doc_result["paragraphs"] = paragraphs
                                doc_result["total_paragraphs"] = len(paragraphs)
                            
                            return doc_result
                except Exception as e:
                    # if logger:
                    #     logger.debug(f"LibreOffice conversion failed: {str(e)}", file_path=filepath)
                    pass
    except Exception as e:
        # if logger:
        #        logger.debug(f"LibreOffice check failed: {str(e)}", file_path=filepath)
        pass
    
    # Try olefile
    try:
        import olefile
        
        if not olefile.isOleFile(filepath):
            return {"error": "File is not a valid OLE2 format .doc file", "filepath": filepath}
        
        # if logger:
        #     logger.debug("Attempting to read .doc with olefile", file_path=filepath)
        
        ole = olefile.OleFileIO(filepath)
        
        if ole.exists('WordDocument'):
            stream = ole.openstream('WordDocument')
            data = stream.read()
            
            text_parts = []
            current_text = []
            for byte in data:
                if 32 <= byte <= 126 or byte in [9, 10, 13]:
                    current_text.append(chr(byte))
                else:
                    if len(current_text) > 3:
                        text_parts.append(''.join(current_text))
                    current_text = []
            
            if current_text and len(current_text) > 3:
                text_parts.append(''.join(current_text))
            
            text = ' '.join(text_parts)
            
            ole.close()
            
            if text.strip():
                doc_result = {
                    "filepath": filepath,
                    "text": text,
                    "method": "olefile",
                    "total_characters": len(text),
                    "total_lines": len(text.splitlines()),
                    "note": "Basic text extraction - formatting may be lost"
                }
                
                paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                if paragraphs:
                    doc_result["paragraphs"] = paragraphs
                    doc_result["total_paragraphs"] = len(paragraphs)
                
                return doc_result
            else:
                ole.close()
                return {"error": "Could not extract text from .doc file", "filepath": filepath}
        else:
            ole.close()
            return {"error": "WordDocument stream not found in .doc file", "filepath": filepath}
            
    except ImportError:
        pass
    except Exception as e:
        # if logger:
        #     logger.debug(f"olefile extraction failed: {str(e)}", file_path=filepath)
        pass
    
    # All methods failed
    error_msg = (
        "Could not read .doc file. Please install one of the following:\n"
        "  - textract: pip install textract (requires antiword or LibreOffice)\n"
        "  - python-docx2txt: pip install docx2txt\n"
        "  - olefile: pip install olefile\n"
        "Or install external tools:\n"
        "  - antiword (command-line tool)\n"
        "  - LibreOffice (soffice command)"
    )
    
    # if logger:
    #     logger.warning("All .doc reading methods failed", file_path=filepath)
    
    return {
        "error": error_msg,
        "filepath": filepath
    }


def read_xlsx_file(filepath):
    """Independent XLSX reader function."""
    try:
        import openpyxl
    except ImportError:
        return {
            "error": "openpyxl not installed. Install with: pip install openpyxl",
            "filepath": filepath
        }
    
    try:
        if not os.path.exists(filepath):
            return {"error": "File not found", "filepath": filepath}
        
        workbook = openpyxl.load_workbook(filepath, data_only=True)
        
        result = {
            "filepath": filepath,
            "sheet_names": workbook.sheetnames,
            "sheets": {},
            "total_sheets": len(workbook.sheetnames)
        }
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = {
                "name": sheet_name,
                "max_row": sheet.max_row,
                "max_column": sheet.max_column,
                "data": []
            }
            
            for row in sheet.iter_rows(values_only=True):
                sheet_data["data"].append(list(row))
            
            result["sheets"][sheet_name] = sheet_data
        
        return result
    except Exception as e:
        return {"error": str(e), "filepath": filepath}


def read_xls_file(filepath):
    """Independent XLS reader function."""
    try:
        import xlrd
    except ImportError:
        return {
            "error": (
                "xlrd not installed (or wrong version). "
                "Install with: pip install 'xlrd==1.2.0'"
            ),
            "filepath": filepath
        }

    try:
        if not os.path.exists(filepath):
            return {"error": "File not found", "filepath": filepath}

        workbook = xlrd.open_workbook(filepath)

        result = {
            "filepath": filepath,
            "sheet_names": workbook.sheet_names(),
            "sheets": {},
            "total_sheets": len(workbook.sheet_names())
        }

        for sheet_name in workbook.sheet_names():
            sheet = workbook.sheet_by_name(sheet_name)
            sheet_data = {
                "name": sheet_name,
                "nrows": sheet.nrows,
                "ncols": sheet.ncols,
                "data": []
            }

            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                sheet_data["data"].append(row_values)

            result["sheets"][sheet_name] = sheet_data

        return result
    except Exception as e:
        return {"error": str(e), "filepath": filepath}


def read_csv_file(filepath, delimiter=',', encoding='utf-8'):
    """Independent CSV reader function."""
    try:
        import csv
        
        if not os.path.exists(filepath):
            return {"error": "File not found", "filepath": filepath}
        
        result = {
            "filepath": filepath,
            "delimiter": delimiter,
            "headers": [],
            "rows": []
        }
        
        encodings = [encoding, 'utf-8', 'latin-1', 'cp1252']
        
        for enc in encodings:
            try:
                with open(filepath, 'r', encoding=enc, newline='') as file:
                    reader = csv.reader(file, delimiter=delimiter)
                    
                    try:
                        result["headers"] = next(reader)
                    except StopIteration:
                        return {"error": "Empty file", "filepath": filepath}
                    
                    for row in reader:
                        result["rows"].append(row)
                
                result["encoding_used"] = enc
                break
            except UnicodeDecodeError:
                if enc == encodings[-1]:
                    raise
                continue
        
        result["row_count"] = len(result["rows"])
        result["column_count"] = len(result["headers"])
        
        return result
    except Exception as e:
        return {"error": str(e), "filepath": filepath}


def read_pptx_file(filepath):
    """Independent PPTX reader function."""
    try:
        from pptx import Presentation
    except ImportError:
        return {
            "error": "python-pptx not installed. Install with: pip install python-pptx",
            "filepath": filepath
        }
    
    try:
        if not os.path.exists(filepath):
            return {"error": "File not found", "filepath": filepath}
        
        prs = Presentation(filepath)
        
        result = {
            "filepath": filepath,
            "slides": [],
            "total_slides": len(prs.slides)
        }
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": slide_idx + 1,
                "texts": []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        slide_data["texts"].append(shape.text)
            
            result["slides"].append(slide_data)
        
        return result
    except Exception as e:
        return {"error": str(e), "filepath": filepath}